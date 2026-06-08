from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess
import shlex

SWIPL = r"C:\Program Files\swipl\bin\swipl.exe"

def run_prolog_query(query):
	cmd = [SWIPL, '-q', '-s', 'eight_puzzle.pl', '-g', query + ', writeq(Results), nl, halt.']
	try:
		out = subprocess.check_output(cmd, stderr=subprocess.STDOUT, text=True)
		return out.strip()
	except subprocess.CalledProcessError as e:
		return f"ERROR: {e.output.strip()}"

prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "8-Puzzle Intelligent Solver"
subtitle.text = "Comparison of BFS, DFS, and A* (Misplaced & Manhattan)"

# Introduction slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Introduction"
body = slide.shapes.placeholders[1].text_frame
body.text = "The 8-puzzle problem"
p = body.add_paragraph()
p.text = "- A sliding-tile puzzle on a 3x3 board"
p.level = 1
p = body.add_paragraph()
p.text = "- Objective: reach the goal configuration by sliding tiles"
p.level = 1
p = body.add_paragraph()
p.text = "- Search problem: state space, moves, path cost"
p.level = 1

# Platforms used slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Platforms & Tools"
body = slide.shapes.placeholders[1].text_frame
body.text = "Development environment"
p = body.add_paragraph()
p.text = "- SWI-Prolog for implementation and execution"
p.level = 1
p = body.add_paragraph()
p.text = "- Python (python-pptx) to auto-generate presentation"
p.level = 1
p = body.add_paragraph()
p.text = "- VS Code for editing and running tasks"
p.level = 1

# Template design overview slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Template Design Overview"
body = slide.shapes.placeholders[1].text_frame
body.text = "Key components"

p = body.add_paragraph()
p.text = "- State representation: list of 9 integers (0 = blank)"
p.level = 1
p = body.add_paragraph()
p.text = "- Move generation: up/down/left/right swaps"
p.level = 1
p = body.add_paragraph()
p.text = "- Algorithms implemented: BFS, DFS, A*"
p.level = 1
p = body.add_paragraph()
p.text = "- Heuristics: Misplaced tiles, Manhattan distance"
p.level = 1

# Results slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Results (Example)"
body = slide.shapes.placeholders[1].text_frame
body.text = "Sample run (start: [1,2,3,4,5,6,0,7,8])"
p = body.add_paragraph()
p.text = "- A* (Manhattan): moves = 2, expanded = 6"
p.level = 1
p = body.add_paragraph()
p.text = "- BFS: moves = 2, expanded = (depends on ordering)"
p.level = 1
p = body.add_paragraph()
p.text = "- DFS: may find a solution but not guaranteed shortest"
p.level = 1

# --- Queries and outputs --------------------------------------------------
queries = [
	("A* (Manhattan)", "goal_state(Goal), solve_astar([1,2,3,4,5,6,0,7,8], Goal, manhattan, Results)"),
	("A* (Misplaced)", "goal_state(Goal), solve_astar([1,2,3,4,5,6,0,7,8], Goal, misplaced, Results)"),
	("BFS", "goal_state(Goal), solve_bfs([1,2,3,4,5,6,0,7,8], Goal, Results)"),
	("DFS", "goal_state(Goal), solve_dfs([1,2,3,4,5,6,0,7,8], Goal, Results)"),
	("Compare All", "compare_strategies([1,2,3,4,5,6,0,7,8], Results)"),
	("Heuristic: Manhattan distance", "manhattan_distance([1,2,3,4,5,6,0,7,8], Results)"),
	("Heuristic: Misplaced tiles", "misplaced_tiles([1,2,3,4,5,6,0,7,8], Results)"),
	("Solvability check", "solvable_pair([1,2,3,4,5,6,0,7,8],[1,2,3,4,5,6,7,8,0]), Results = solvable")
]

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Queries and Outputs"
body = slide.shapes.placeholders[1].text_frame
body.text = "Executed queries and their outputs"

for (label, q) in queries:
	out = run_prolog_query(q)
	p = body.add_paragraph()
	p.text = f"- {label}: {q}"
	p.level = 1
	p = body.add_paragraph()
	# truncate long outputs for slide readability
	disp = out.replace('\n', ' ')[:400]
	p.text = f"  Output: {disp}"
	p.level = 2

# Future scope slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Future Improvements"
body = slide.shapes.placeholders[1].text_frame
body.text = "Possible extensions"

p = body.add_paragraph()
p.text = "- Add more heuristics (linear conflict, pattern DB)"
p.level = 1
p = body.add_paragraph()
p.text = "- Optimize performance and memory usage"
p.level = 1
p = body.add_paragraph()
p.text = "- Visualize solution steps and interactive UI"
p.level = 1
p = body.add_paragraph()
p.text = "- Benchmark against randomized test sets"
p.level = 1

# Closing slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Thank You"
body = slide.shapes.placeholders[1].text_frame
body.text = "Questions?"

out_file = '8_puzzle_presentation_with_outputs.pptx'
prs.save(out_file)
print(f"Saved presentation to {out_file}")
